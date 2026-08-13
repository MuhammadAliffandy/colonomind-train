import os
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "ColonoMind System Architecture & Results"
slide.placeholders[1].text = "Adaptive Ensemble & Unified Dataset Backbone\nAugust 2026"

blank_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = "System Architecture Overview"
try:
    slide.shapes.add_picture("/Users/aliffandy/.gemini/antigravity-ide/brain/cf8e5810-1798-4ed3-a684-5aa8ab0bc120/colonomind_architecture_bw_final_1786459359603.png", Inches(1), Inches(1.5), height=Inches(5))
except:
    pass

bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Unified Dataset Backbone"
tf = slide.placeholders[1].text_frame
tf.text = "The system is now powered by a 'Unified' training scenario."
p = tf.add_paragraph(); p.text = "Combines 3 Massive Datasets:"; p.level = 1
p = tf.add_paragraph(); p.text = "TMC-UCM, NTUH, and LIMUC"; p.level = 2
p = tf.add_paragraph(); p.text = "Advantages:"; p.level = 1
p = tf.add_paragraph(); p.text = "Increases generalization across different patient demographics and camera hardware."; p.level = 2

slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "Smart Ensemble & Decision Logic"
tf = slide.placeholders[1].text_frame
tf.text = "Multi-Model Fusion:"
p = tf.add_paragraph(); p.text = "Uses 5 distinct deep learning models (ResNet, DenseNet, EfficientNet, ConvNeXt, ViT) alongside LightGBM."; p.level = 1
p = tf.add_paragraph(); p.text = "Majority Voting & Average Probability:"; p.level = 0
p = tf.add_paragraph(); p.text = "Final consensus requires at least 3 out of 5 models to agree on the severity score."; p.level = 1
p = tf.add_paragraph(); p.text = "Fallback Most-Severe MES Rules:"; p.level = 0
p = tf.add_paragraph(); p.text = "When models deeply disagree (Consensus < Threshold), the system prioritizes patient safety by highlighting the Most-Severe prediction."; p.level = 1

metrics = {
    "Intra TMC-UCM": [
        {"Model": "EfficientNet-B4", "Base Acc": "79.44%", "Hyb Acc": "78.47%", "Prec": "77.20%", "Recall": "76.07%", "F1": "76.32%", "QWK": "0.9018"},
        {"Model": "ResNet-50",       "Base Acc": "81.79%", "Hyb Acc": "80.57%", "Prec": "79.82%", "Recall": "79.13%", "F1": "79.42%", "QWK": "0.9051"},
        {"Model": "ViT-B-16",        "Base Acc": "79.84%", "Hyb Acc": "75.02%", "Prec": "74.58%", "Recall": "74.15%", "F1": "74.32%", "QWK": "0.8476"},
        {"Model": "ConvNeXt-Tiny",   "Base Acc": "80.35%", "Hyb Acc": "78.53%", "Prec": "77.89%", "Recall": "76.66%", "F1": "77.19%", "QWK": "0.8899"},
        {"Model": "DenseNet-121",    "Base Acc": "79.63%", "Hyb Acc": "75.86%", "Prec": "74.23%", "Recall": "74.24%", "F1": "74.23%", "QWK": "0.8702"},
    ],
    "Intra Unified": [
        {"Model": "ConvNeXt-Tiny",   "Base Acc": "79.06%", "Hyb Acc": "79.50%", "Prec": "78.19%", "Recall": "75.65%", "F1": "76.63%", "QWK": "0.9000"},
        {"Model": "ResNet-50",       "Base Acc": "77.33%", "Hyb Acc": "78.34%", "Prec": "76.15%", "Recall": "74.32%", "F1": "75.16%", "QWK": "0.8950"},
        {"Model": "ViT-B-16",        "Base Acc": "77.82%", "Hyb Acc": "78.22%", "Prec": "75.22%", "Recall": "75.69%", "F1": "75.41%", "QWK": "0.8939"},
        {"Model": "EfficientNet-B4", "Base Acc": "76.41%", "Hyb Acc": "76.73%", "Prec": "73.52%", "Recall": "73.51%", "F1": "73.43%", "QWK": "0.8871"},
        {"Model": "DenseNet-121",    "Base Acc": "76.46%", "Hyb Acc": "76.39%", "Prec": "75.12%", "Recall": "71.74%", "F1": "73.22%", "QWK": "0.8751"},
    ],
    "Intra LIMUC": [
        {"Model": "EfficientNet-B4", "Base Acc": "74.43%", "Hyb Acc": "69.69%", "Prec": "57.87%", "Recall": "57.02%", "F1": "57.22%", "QWK": "0.7070"},
        {"Model": "ResNet-50",       "Base Acc": "75.91%", "Hyb Acc": "74.67%", "Prec": "66.85%", "Recall": "65.68%", "F1": "66.14%", "QWK": "0.8179"},
        {"Model": "ViT-B-16",        "Base Acc": "74.37%", "Hyb Acc": "68.62%", "Prec": "58.91%", "Recall": "58.72%", "F1": "58.59%", "QWK": "0.7058"},
        {"Model": "ConvNeXt-Tiny",   "Base Acc": "76.86%", "Hyb Acc": "71.41%", "Prec": "61.15%", "Recall": "61.27%", "F1": "61.09%", "QWK": "0.7353"},
        {"Model": "DenseNet-121",    "Base Acc": "75.74%", "Hyb Acc": "71.05%", "Prec": "62.51%", "Recall": "61.78%", "F1": "62.08%", "QWK": "0.7182"},
    ],
    "Intra NTUH": [
        {"Model": "EfficientNet-B4", "Base Acc": "62.81%", "Hyb Acc": "64.32%", "Prec": "61.22%", "Recall": "61.98%", "F1": "61.06%", "QWK": "0.7238"},
        {"Model": "ResNet-50",       "Base Acc": "67.83%", "Hyb Acc": "66.33%", "Prec": "62.26%", "Recall": "61.35%", "F1": "60.55%", "QWK": "0.7421"},
        {"Model": "ViT-B-16",        "Base Acc": "65.82%", "Hyb Acc": "67.33%", "Prec": "64.19%", "Recall": "63.47%", "F1": "63.11%", "QWK": "0.7424"},
        {"Model": "ConvNeXt-Tiny",   "Base Acc": "68.84%", "Hyb Acc": "66.33%", "Prec": "63.60%", "Recall": "64.35%", "F1": "63.68%", "QWK": "0.7446"},
        {"Model": "DenseNet-121",    "Base Acc": "69.84%", "Hyb Acc": "68.34%", "Prec": "64.63%", "Recall": "65.35%", "F1": "64.47%", "QWK": "0.7892"},
    ],
}

headers = ["Model", "Base Acc", "Hyb Acc", "Prec", "Recall", "F1", "QWK"]

for dataset_name, data_list in metrics.items():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Experimental Results: {dataset_name}"

    rows = len(data_list) + 1
    cols = len(headers)
    
    # Table dimensions
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.4 * rows)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    for i, h in enumerate(headers): table.cell(0, i).text = h
        
    for r, data in enumerate(data_list):
        table.cell(r+1, 0).text = data["Model"]
        table.cell(r+1, 1).text = data["Base Acc"]
        table.cell(r+1, 2).text = data["Hyb Acc"]
        table.cell(r+1, 3).text = data["Prec"]
        table.cell(r+1, 4).text = data["Recall"]
        table.cell(r+1, 5).text = data["F1"]
        table.cell(r+1, 6).text = data["QWK"]

    for row in table.rows:
        row.height = Pt(18)
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)

prs.save("ColonoMind_System_Architecture_and_Results_FINAL.pptx")
print("Saved final PPT.")
