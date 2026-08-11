import os
import json
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ColonoMind System Architecture & Results"
subtitle.text = "Adaptive Ensemble & Unified Dataset Backbone\nAugust 2026"

# Slide 2: Architecture Diagram
blank_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(blank_slide_layout)
title = slide.shapes.title
title.text = "System Architecture Overview"

img_path = "assets/architecture.png"
# Add image (centered)
try:
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), height=Inches(5))
except Exception as e:
    print(f"Image not found or error loading: {e}")

# Slide 3: The Unified Backbone
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Unified Dataset Backbone"
tf = body_shape.text_frame
tf.text = "The system is now powered by a 'Unified' training scenario."
p = tf.add_paragraph()
p.text = "Combines 3 Massive Datasets:"
p.level = 1
p = tf.add_paragraph()
p.text = "TMC-UCM, NTUH, and LIMUC"
p.level = 2
p = tf.add_paragraph()
p.text = "Advantages:"
p.level = 1
p = tf.add_paragraph()
p.text = "Increases generalization across different patient demographics and camera hardware."
p.level = 2
p = tf.add_paragraph()
p.text = "Overcomes the limitations of cross-domain evaluation drops."
p.level = 2

# Slide 4: Adaptive Ensemble & Majority Voting
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Smart Ensemble & Decision Logic"
tf = body_shape.text_frame
tf.text = "Multi-Model Fusion:"
p = tf.add_paragraph()
p.text = "Uses 5 distinct deep learning models (ResNet, DenseNet, EfficientNet, ConvNeXt, ViT) alongside LightGBM."
p.level = 1
p = tf.add_paragraph()
p.text = "Majority Voting & Average Probability:"
p.level = 0
p = tf.add_paragraph()
p.text = "Final consensus requires at least 3 out of 5 models to agree on the severity score."
p.level = 1
p = tf.add_paragraph()
p.text = "Fallback Most-Severe MES Rules:"
p.level = 0
p = tf.add_paragraph()
p.text = "When models deeply disagree (Consensus < Threshold), the system prioritizes patient safety by highlighting the Most-Severe prediction."
p.level = 1

# Slide 5: Current Experimental Results
# Parse the results from the Result folder dynamically
result_dir = "Result"
metrics = []
if os.path.exists(result_dir):
    json_files = glob.glob(os.path.join(result_dir, "*", "*", "*_metrics.json"))
    for jf in json_files:
        try:
            with open(jf, 'r') as f:
                data = json.load(f)
                
            parts = jf.split(os.sep)
            dataset = parts[-3]  # e.g. Intra_Unified
            model = parts[-2].replace("_Experiment", "") # e.g. ResNet-50
            
            acc = data.get("test_accuracy", 0) * 100
            f1 = data.get("test_f1", 0) * 100
            auc = data.get("test_auc", 0)
            
            metrics.append({
                "Dataset": dataset,
                "Model": model,
                "Accuracy": f"{acc:.2f}%",
                "F1 Score": f"{f1:.2f}%",
                "AUC": f"{auc:.4f}"
            })
        except Exception as e:
            print(f"Error parsing {jf}: {e}")
else:
    print(f"Result directory '{result_dir}' not found. Cannot populate results table.")

slide_layout = prs.slide_layouts[5] # blank slide for table
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Current Experimental Results"

if not metrics:
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    tf = txBox.text_frame
    tf.text = "No results found in the Result/ directory yet."
else:
    # Sort metrics logically
    metrics = sorted(metrics, key=lambda x: (x["Dataset"], x["Model"]))
    
    rows = len(metrics) + 1
    cols = 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.4 * rows)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Headers
    headers = ["Dataset", "Model", "Accuracy", "F1 Score", "AUC"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    # Rows
    for r, row_data in enumerate(metrics):
        table.cell(r+1, 0).text = row_data["Dataset"]
        table.cell(r+1, 1).text = row_data["Model"]
        table.cell(r+1, 2).text = row_data["Accuracy"]
        table.cell(r+1, 3).text = row_data["F1 Score"]
        table.cell(r+1, 4).text = row_data["AUC"]
        
    # Adjust font sizes
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)

# Save the presentation
output_name = "ColonoMind_System_Architecture_and_Results.pptx"
prs.save(output_name)
print(f"Presentation saved as {output_name}")
